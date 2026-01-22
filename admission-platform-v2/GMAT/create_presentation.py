from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RGBColor(41, 98, 255)      # Blue
SECONDARY_COLOR = RGBColor(0, 200, 150)    # Teal/Green
ACCENT_COLOR = RGBColor(255, 107, 107)     # Coral/Red
DARK_COLOR = RGBColor(30, 30, 60)          # Dark blue
LIGHT_BG = RGBColor(245, 247, 250)         # Light gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with gradient-style background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()

    # Decorative element
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.5), prs.slide_width, Inches(2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(31, 78, 200)
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_COLOR
    bg.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2), Inches(2.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY_COLOR
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_first=False):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(60, 60, 80)
        p.space_after = Pt(12)

        if highlight_first and i == 0:
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR

    return slide

def add_stats_slide(prs, title, stats):
    """Add a slide with statistics boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # Stats boxes
    colors = [PRIMARY_COLOR, SECONDARY_COLOR, ACCENT_COLOR, RGBColor(150, 100, 200)]
    box_width = Inches(2.8)
    box_height = Inches(2.2)
    start_x = Inches(0.7)
    gap = Inches(0.3)
    y_pos = Inches(2.2)

    for i, (value, label) in enumerate(stats[:4]):
        x_pos = start_x + i * (box_width + gap)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = colors[i % len(colors)]
        box.line.width = Pt(3)

        # Top accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, box_width, Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = colors[i % len(colors)]
        accent.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.4), box_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(1.4), box_width - Inches(0.2), Inches(0.8))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 100, 120)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # Left column
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
    left_box.line.color.rgb = PRIMARY_COLOR
    left_box.line.width = Pt(2)

    left_title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(5.4), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    left_content = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(5.2), Inches(4.2))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 80)
        p.space_after = Pt(8)

    # Right column
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 255, 250)
    right_box.line.color.rgb = SECONDARY_COLOR
    right_box.line.width = Pt(2)

    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.4), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    right_content = slide.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.2))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(60, 60, 80)
        p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(11.5)
    table_height = Inches(0.5) * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.9), Inches(1.5), table_width, table_height).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(14)
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(13)
            para.font.color.rgb = DARK_COLOR
            para.alignment = PP_ALIGN.CENTER

    return slide

def add_timeline_slide(prs, title, weeks):
    """Add a timeline slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_COLOR
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(10), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(200, 200, 210)
    line.line.fill.background()

    colors = [PRIMARY_COLOR, SECONDARY_COLOR, RGBColor(150, 100, 200), ACCENT_COLOR]
    box_width = Inches(2.3)
    start_x = Inches(0.8)
    gap = Inches(0.35)

    for i, (week_title, tasks) in enumerate(weeks):
        x_pos = start_x + i * (box_width + gap)

        # Circle on timeline
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(1), Inches(2.35), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()

        # Week box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3), box_width, Inches(4))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = colors[i]
        box.line.width = Pt(2)

        # Week title
        week_box = slide.shapes.add_textbox(x_pos, Inches(3.1), box_width, Inches(0.5))
        tf = week_box.text_frame
        p = tf.paragraphs[0]
        p.text = week_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER

        # Tasks
        tasks_box = slide.shapes.add_textbox(x_pos + Inches(0.15), Inches(3.6), box_width - Inches(0.3), Inches(3.2))
        tf = tasks_box.text_frame
        tf.word_wrap = True
        for j, task in enumerate(tasks):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + task
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(80, 80, 100)
            p.space_after = Pt(4)

    return slide

def add_closing_slide(prs, title, subtitle, items):
    """Add a closing/summary slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_COLOR
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Items
    items_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.333), Inches(3.5))
    tf = items_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "✓  " + item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT

    return slide

# ============================================
# CREATE THE PRESENTATION
# ============================================

# Slide 1: Title
add_title_slide(prs,
    "GMAT Question Generation\nvia AI Fine-Tuning",
    "Strategic Proposal for UpToTen Test Platform"
)

# Slide 2: The Challenge - Stats
add_stats_slide(prs, "The Challenge: Question Gap Analysis", [
    ("1,342", "Current Questions\nAvailable"),
    ("1,172", "Questions Needed\nPer Student/Cycle"),
    ("5,000+", "Additional Questions\nNeeded for Scale"),
    ("4 weeks", "Target Timeline\nto Fill Gap"),
])

# Slide 3: Current Distribution
add_table_slide(prs, "Current Question Distribution",
    ["Section", "Available", "Per Cycle Need", "Status"],
    [
        ["Quantitative Reasoning", "582", "469", "✓ 1 cycle"],
        ["Verbal Reasoning", "383", "311", "✓ 1 cycle"],
        ["Data Insights", "377", "392", "⚠ Slight gap"],
        ["TOTAL", "1,342", "1,172", "Gap grows with scale"],
    ]
)

# Slide 4: Section divider
add_section_slide(prs, "The Proposed Solution")

# Slide 5: What is Fine-Tuning
add_content_slide(prs, "What is AI Fine-Tuning?", [
    "Adapts a pre-trained AI model (GPT-4o) to your specific use case",
    "Trains on examples of desired output format and style",
    "Unlike prompting, it modifies the model's behavior permanently",
    "Result: Consistent, high-quality GMAT-style questions on demand",
    "One-time training → Unlimited generation capability",
])

# Slide 6: Comparison
add_comparison_slide(prs, "Why Fine-Tuning?",
    "Fine-Tuning Approach",
    [
        "Low cost: ~$50 for training + generation",
        "High consistency: Same format every time",
        "Scalable: Generate thousands on demand",
        "Fast after setup: Seconds per question",
        "Token-efficient: No examples needed in prompts",
    ],
    "Alternatives",
    [
        "Manual creation: $25,000+ for 2,500 questions",
        "Few-shot prompting: Higher per-question cost",
        "Purchase questions: $5,000-15,000 licensing",
        "Partner with test prep: Revenue share + dependency",
        "All alternatives: Slower or more expensive",
    ]
)

# Slide 7: How It Works
add_content_slide(prs, "How It Works: The Pipeline", [
    "STEP 1: Select 150-300 of our best existing questions for training",
    "STEP 2: Convert to OpenAI's JSONL training format",
    "STEP 3: Fine-tune GPT-4o model (2-4 hours, ~$10)",
    "STEP 4: Generate new questions via API (~$8 per 1,000)",
    "STEP 5: Expert review and validation (in-house GMAT experts)",
    "STEP 6: Import approved questions to platform",
])

# Slide 8: Section divider
add_section_slide(prs, "Implementation Plan")

# Slide 9: Timeline
add_timeline_slide(prs, "4-Week Implementation Timeline", [
    ("Week 1", [
        "Select 150 training questions",
        "Build export scripts",
        "Set up OpenAI API",
        "Run first fine-tuning",
        "Generate 50 pilot questions",
    ]),
    ("Week 2", [
        "Expert reviews pilot",
        "Iterate on training",
        "Generate batch 1 (500)",
        "Begin parallel review",
        "Quality assessment",
    ]),
    ("Week 3", [
        "Scale generation (2,000+)",
        "Parallel expert review",
        "Fix/reject as needed",
        "Quality tracking",
        "Refinement cycle",
    ]),
    ("Week 4", [
        "Convert to platform format",
        "Database import",
        "Platform testing",
        "Final QA",
        "Launch!",
    ]),
])

# Slide 10: Cost Analysis
add_stats_slide(prs, "Cost Analysis", [
    ("~$50", "AI Costs\n(Training + Generation)"),
    ("65 hrs", "Total Human Effort\n(Dev + Expert)"),
    ("$3,500", "Total Estimated Cost\n(at $50/hour)"),
    ("86%", "Cost Savings vs.\nManual Creation"),
])

# Slide 11: Cost Comparison Table
add_table_slide(prs, "Cost Comparison: AI vs. Alternatives",
    ["Approach", "Questions", "Time", "Cost"],
    [
        ["Manual Creation", "2,500", "500+ hours", "$25,000+"],
        ["Purchase/License", "2,500", "2-4 weeks", "$5,000-15,000"],
        ["AI + Expert Review", "2,500", "65 hours", "~$3,500"],
        ["SAVINGS with AI", "—", "87% less time", "86% less cost"],
    ]
)

# Slide 12: Section divider
add_section_slide(prs, "Quality Assurance")

# Slide 13: QA Process
add_comparison_slide(prs, "Quality Assurance Process",
    "Automated Checks",
    [
        "Format validation: Correct JSON structure",
        "Field verification: All required fields present",
        "Answer options: Exactly 5 choices (a-e)",
        "Mathematical validation: Verify calculations",
        "Duplication check: Compare to existing questions",
    ],
    "Expert Review",
    [
        "Mathematical/logical correctness",
        "Clear, unambiguous wording",
        "Appropriate difficulty level",
        "Realistic GMAT-style scenarios",
        "Quality explanation verification",
        "Target: 100-150 questions/expert/day",
    ]
)

# Slide 14: Success Metrics
add_content_slide(prs, "Success Metrics & Go/No-Go Criteria", [
    "QUALITY: >80% of generated questions pass expert review on first pass",
    "ACCURACY: <5% of questions contain mathematical/logical errors",
    "SPEED: 2,500+ validated questions delivered within 4 weeks",
    "COVERAGE: Balanced distribution across all topics and difficulty levels",
    "",
    "GO/NO-GO CHECKPOINTS:",
    "   Day 5: Pilot review - >70% pass rate to continue",
    "   Week 2: <10% error rate to scale up",
    "   Week 4: 2,000+ validated questions to launch",
])

# Slide 15: Section divider
add_section_slide(prs, "Risk Assessment")

# Slide 16: Risks
add_table_slide(prs, "Risk Assessment & Mitigations",
    ["Risk", "Impact", "Mitigation", "Residual"],
    [
        ["Quality Issues", "High", "100% expert review, pilot first", "Low"],
        ["Copyright/Legal", "Medium", "Train on format, not content; legal review", "Low-Med"],
        ["Difficulty Calibration", "Medium", "Include labels, track performance", "Low"],
        ["Timeline Slippage", "Low", "Buffer in schedule, parallel work", "Low"],
    ]
)

# Slide 17: Section divider
add_section_slide(prs, "Decision & Next Steps")

# Slide 18: Required Resources
add_content_slide(prs, "Resources Required", [
    "BUDGET:",
    "   • ~$50 for AI costs (training + generation)",
    "   • Expert time: 25-35 hours over 4 weeks",
    "   • Developer time: 20-30 hours over 4 weeks",
    "",
    "PERSONNEL:",
    "   • 1 Developer (scripts, integration, automation)",
    "   • 1-2 GMAT Experts (training selection, validation)",
    "",
    "TECHNICAL:",
    "   • OpenAI API account with payment method",
    "   • Access to existing question database",
])

# Slide 19: Deliverables
add_content_slide(prs, "Deliverables", [
    "JSONL export script - Convert existing questions to training format",
    "Fine-tuned GPT-4o model - Specialized for GMAT question generation",
    "Generation scripts - Automated question creation with validation",
    "Expert review system - Tracking spreadsheet for QA process",
    "2,500+ validated questions - Imported and ready in platform",
    "Process documentation - Guide for future question generation",
])

# Slide 20: Closing
add_closing_slide(prs,
    "Recommendation: Proceed",
    "AI Fine-Tuning offers the best balance of speed, cost, and quality",
    [
        "Fill the 5,000+ question gap in just 4 weeks",
        "Save 86% compared to manual creation",
        "In-house experts ensure quality control",
        "Scalable solution for future growth",
        "Low-risk pilot approach (50 questions first)",
    ]
)

# Slide 21: Final slide
add_title_slide(prs,
    "Questions?",
    "Ready to discuss next steps"
)

# Save the presentation
output_path = r"c:\Users\kappa\Desktop\UpToTen\test-ammissione\admission-platform-v2\GMAT\GMAT_AI_Question_Generation_Proposal.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
